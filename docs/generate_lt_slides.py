"""LT(ライトニングトーク)用スライド生成スクリプト — 技術勉強会向け10枚構成."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Catppuccin Mocha カラーパレット
BG_PRIMARY = RGBColor(0x1E, 0x1E, 0x2E)
BG_SURFACE = RGBColor(0x31, 0x32, 0x44)
TEXT_PRIMARY = RGBColor(0xCD, 0xD6, 0xF4)
TEXT_SECONDARY = RGBColor(0xA6, 0xAD, 0xC8)
TEXT_MUTED = RGBColor(0x6C, 0x70, 0x86)
ACCENT = RGBColor(0x89, 0xB4, 0xFA)
SUCCESS = RGBColor(0xA6, 0xE3, 0xA1)
WARNING = RGBColor(0xF9, 0xE2, 0xAF)
ERROR = RGBColor(0xF3, 0x8B, 0xA8)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_bg(slide, color=BG_PRIMARY):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def txt(slide, left, top, width, height, text, size=18,
        color=TEXT_PRIMARY, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Yu Gothic UI'
    p.alignment = align
    return box


def rect(slide, left, top, width, height, fill, text="",
         size=14, color=TEXT_PRIMARY):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(size)
        tf.paragraphs[0].font.color.rgb = color
        tf.paragraphs[0].font.name = 'Yu Gothic UI'
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def multi(slide, left, top, width, height, lines, default_size=18):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        if isinstance(ln, str):
            ln = {'text': ln}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln.get('text', '')
        p.font.size = Pt(ln.get('size', default_size))
        p.font.color.rgb = ln.get('color', TEXT_PRIMARY)
        p.font.bold = ln.get('bold', False)
        p.font.name = 'Yu Gothic UI'
        p.space_after = Pt(ln.get('space', 6))
        p.alignment = ln.get('align', PP_ALIGN.LEFT)
    return box


def circle_num(slide, x, y, num, color):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.9), Inches(0.9))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    tf = c.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.color.rgb = BG_PRIMARY
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = 'Yu Gothic UI'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def table_row(slide, y, cols, sizes, colors, bolds, bg=None):
    if bg:
        rect(slide, Inches(0.6), y - Inches(0.05), Inches(12), Inches(0.45), bg)
    for col_text, size, color, bold, x in zip(
            cols, sizes, colors, bolds,
            [Inches(0.8), Inches(3.6), Inches(6.4), Inches(9.5)]):
        txt(slide, x, y, Inches(2.8), Inches(0.4), col_text,
            size=size, color=color, bold=bold)


def create():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    qr = os.path.join(os.path.dirname(__file__), 'LP用QRコード.png')

    # ================================================================
    # 1: タイトル
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    txt(s, Inches(0.8), Inches(0.8), Inches(11.5), Inches(1.2),
        'スクラッチ開発未経験エンジニアが\nAI駆動開発で3週間でアプリを作った話',
        size=38, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(1.5), Inches(2.8), Inches(10), Inches(0.6),
        '— ユメログ (YumeLog) 開発記 —',
        size=22, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
    if os.path.exists(qr):
        s.shapes.add_picture(qr, Inches(5.4), Inches(3.8), Inches(2.5), Inches(2.5))
    txt(s, Inches(3), Inches(6.5), Inches(7), Inches(0.4),
        'QRコードは最後にもう一度出します。まずは話を聞いてください。',
        size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # ================================================================
    # 2: 自己紹介
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    txt(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
        '自己紹介', size=36, color=ACCENT, bold=True)

    rect(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(5.3), BG_SURFACE)
    multi(s, Inches(1.2), Inches(1.6), Inches(4.8), Inches(5.0), [
        {'text': 'プロフィール', 'size': 22, 'color': ACCENT, 'bold': True, 'space': 16},
        {'text': '27歳 / IT業界 5年目', 'size': 20, 'bold': True, 'space': 20},
        {'text': 'IT業界経験', 'size': 13, 'color': TEXT_MUTED, 'space': 2},
        {'text': 'パッケージ開発（約3年）', 'size': 16, 'space': 4},
        {'text': 'ローコード開発（約3年）', 'size': 16, 'space': 16},
        {'text': 'スクラッチ開発経験', 'size': 13, 'color': TEXT_MUTED, 'space': 2},
        {'text': 'なし（ユメログが初）', 'size': 16, 'color': WARNING, 'bold': True, 'space': 16},
        {'text': 'Flutter / Dart 経験', 'size': 13, 'color': TEXT_MUTED, 'space': 2},
        {'text': 'なし（ユメログが初）', 'size': 16, 'color': WARNING, 'bold': True},
    ])

    rect(s, Inches(6.8), Inches(1.4), Inches(5.7), Inches(5.3), BG_SURFACE)
    multi(s, Inches(7.2), Inches(1.6), Inches(5.0), Inches(5.0), [
        {'text': '今日のゴール', 'size': 22, 'color': ACCENT, 'bold': True, 'space': 24},
        {'text': '', 'size': 8, 'space': 4},
        {'text': 'AI駆動開発で', 'size': 20, 'color': TEXT_SECONDARY, 'space': 4},
        {'text': '何がどれだけ変わるのか', 'size': 24, 'color': TEXT_PRIMARY, 'bold': True, 'space': 4},
        {'text': 'を数字で示す', 'size': 20, 'color': TEXT_SECONDARY, 'space': 28},
        {'text': '', 'size': 8, 'space': 4},
        {'text': 'そして、その成果物を', 'size': 18, 'color': TEXT_MUTED, 'space': 4},
        {'text': '皆さんの目で確かめてもらう', 'size': 18, 'color': SUCCESS, 'bold': True},
    ])

    # ================================================================
    # 3: Before — 従来開発の見積もり
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    txt(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
        'Before: もし従来の開発で作ったら？', size=36, color=ACCENT, bold=True)

    rect(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.8), BG_SURFACE)
    multi(s, Inches(1.2), Inches(1.7), Inches(10.8), Inches(1.5), [
        {'text': '今回作ったアプリの規模', 'size': 20, 'color': ACCENT, 'bold': True, 'space': 12},
        {'text': 'コード約48,000行  /  テスト813件  /  32機能  /  8画面  /  6テーブル',
         'size': 18, 'space': 4},
    ])

    txt(s, Inches(0.8), Inches(3.8), Inches(11), Inches(0.5),
        'IPA ソフトウェア開発データ白書 + COCOMO II モデルで試算すると…',
        size=16, color=TEXT_MUTED)

    items = [
        ('フレームワーク学習', '1〜2ヶ月', 'Flutter / Dart / Drift / Riverpod を学ぶ'),
        ('設計・開発・テスト', '4〜7ヶ月', '生産性 50〜100行/人日（IPA中央値）'),
        ('合計', '6〜9ヶ月', '1人で新規スクラッチ開発の場合'),
    ]
    for i, (label, period, note) in enumerate(items):
        y = Inches(4.4 + i * 0.9)
        bg = BG_SURFACE if i % 2 == 0 else None
        if bg:
            rect(s, Inches(0.8), y - Inches(0.05), Inches(11.5), Inches(0.8), bg)
        c = WARNING if i == 2 else TEXT_PRIMARY
        b = i == 2
        txt(s, Inches(1.2), y + Inches(0.1), Inches(3), Inches(0.5),
            label, size=18, color=c, bold=b)
        txt(s, Inches(4.5), y + Inches(0.1), Inches(2.5), Inches(0.5),
            period, size=20, color=ERROR if i == 2 else ACCENT, bold=True)
        txt(s, Inches(7.5), y + Inches(0.1), Inches(4.5), Inches(0.5),
            note, size=14, color=TEXT_MUTED)

    # ================================================================
    # 4: After — AI駆動開発の実績 ★
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    txt(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
        'After: AI駆動開発の実績', size=36, color=ACCENT, bold=True)

    # 上部カード
    rect(s, Inches(0.8), Inches(1.4), Inches(3.6), Inches(2.0), BG_SURFACE)
    multi(s, Inches(1.2), Inches(1.5), Inches(3.0), Inches(1.8), [
        {'text': '開発期間', 'size': 14, 'color': TEXT_MUTED, 'space': 4},
        {'text': '21日', 'size': 42, 'color': SUCCESS, 'bold': True, 'space': 4, 'align': PP_ALIGN.CENTER},
        {'text': '（従来: 6〜9ヶ月）', 'size': 14, 'color': TEXT_MUTED, 'align': PP_ALIGN.CENTER},
    ])
    rect(s, Inches(4.8), Inches(1.4), Inches(3.6), Inches(2.0), BG_SURFACE)
    multi(s, Inches(5.2), Inches(1.5), Inches(3.0), Inches(1.8), [
        {'text': '1日あたり生産量', 'size': 14, 'color': TEXT_MUTED, 'space': 4},
        {'text': '2,281行', 'size': 42, 'color': SUCCESS, 'bold': True, 'space': 4, 'align': PP_ALIGN.CENTER},
        {'text': '（従来: 50〜100行）', 'size': 14, 'color': TEXT_MUTED, 'align': PP_ALIGN.CENTER},
    ])
    rect(s, Inches(8.8), Inches(1.4), Inches(3.7), Inches(2.0), BG_SURFACE)
    multi(s, Inches(9.2), Inches(1.5), Inches(3.0), Inches(1.8), [
        {'text': '学習コスト', 'size': 14, 'color': TEXT_MUTED, 'space': 4},
        {'text': '0日', 'size': 42, 'color': SUCCESS, 'bold': True, 'space': 4, 'align': PP_ALIGN.CENTER},
        {'text': '（従来: 1〜2ヶ月）', 'size': 14, 'color': TEXT_MUTED, 'align': PP_ALIGN.CENTER},
    ])

    # 比較表
    txt(s, Inches(0.8), Inches(3.8), Inches(6), Inches(0.5),
        '従来開発との比較', size=22, color=WARNING, bold=True)
    header_y = Inches(4.3)
    for label, x in [('指標', 0.8), ('従来開発', 3.6), ('AI駆動開発', 6.4), ('効果', 9.5)]:
        txt(s, Inches(x), header_y, Inches(2.8), Inches(0.4),
            label, size=14, color=TEXT_MUTED, bold=True)

    metrics = [
        ('開発期間', '6〜9ヶ月', '21日', '6〜9倍 短縮'),
        ('コード生産速度', '50〜100行/日', '2,281行/日', '23〜46倍'),
        ('テスト密度', '0〜3件/KLOC', '16.9件/KLOC', '6倍以上'),
        ('テストケース', '0〜50件', '813件', '自動生成'),
    ]
    for i, (a, b, c, d) in enumerate(metrics):
        y = Inches(4.8 + i * 0.55)
        bg = BG_SURFACE if i % 2 == 0 else None
        if bg:
            rect(s, Inches(0.6), y - Inches(0.05), Inches(12), Inches(0.48), bg)
        txt(s, Inches(0.8), y, Inches(2.8), Inches(0.4), a, size=16)
        txt(s, Inches(3.6), y, Inches(2.8), Inches(0.4), b, size=16, color=TEXT_SECONDARY)
        txt(s, Inches(6.4), y, Inches(2.8), Inches(0.4), c, size=16, color=ACCENT, bold=True)
        txt(s, Inches(9.5), y, Inches(3), Inches(0.4), d, size=16, color=SUCCESS, bold=True)

    # ================================================================
    # 5: ライフステージ別比較 ★
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    txt(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
        '工程別の比較 — テスト工程が消滅', size=36, color=ACCENT, bold=True)

    # ヘッダー
    for label, x in [('工程', 0.8), ('AI駆動', 3.2), ('従来開発', 5.5), ('短縮率', 8.0), ('AIの主な貢献', 9.8)]:
        txt(s, Inches(x), Inches(1.3), Inches(2), Inches(0.4),
            label, size=14, color=TEXT_MUTED, bold=True)

    stages = [
        ('企画', '1日', '1〜2週間', '70〜85%', '機能案ブレスト'),
        ('要件定義', '1日', '2〜4週間', '80〜90%', '要件書の自動生成'),
        ('仕様検討', '2日', '3〜6週間', '80〜90%', '画面仕様・DB設計'),
        ('設計', '1日', '2〜4週間', '80〜90%', 'アーキテクチャ選定'),
        ('開発', '14日', '3〜6ヶ月', '85〜95%', 'コード自動生成'),
        ('テスト', '並行', '1〜2ヶ月', '90〜95%', '813件自動生成+CI/CD'),
        ('合計', '約21日', '7〜14ヶ月', '約90%', ''),
    ]
    for i, (stage, ai, trad, rate, contrib) in enumerate(stages):
        y = Inches(1.8 + i * 0.7)
        is_total = i == 6
        is_test = i == 5
        bg = BG_SURFACE if i % 2 == 0 else None
        if is_total:
            bg = ACCENT
        if bg:
            tc = BG_PRIMARY if is_total else TEXT_PRIMARY
            rect(s, Inches(0.6), y - Inches(0.05), Inches(12), Inches(0.58), bg)
        sc = BG_PRIMARY if is_total else (WARNING if is_test else TEXT_PRIMARY)
        sb = is_total or is_test
        txt(s, Inches(0.8), y + Inches(0.05), Inches(2), Inches(0.4),
            stage, size=17, color=sc, bold=sb)
        txt(s, Inches(3.2), y + Inches(0.05), Inches(2), Inches(0.4),
            ai, size=17, color=BG_PRIMARY if is_total else ACCENT, bold=True)
        txt(s, Inches(5.5), y + Inches(0.05), Inches(2.2), Inches(0.4),
            trad, size=17, color=BG_PRIMARY if is_total else TEXT_SECONDARY)
        txt(s, Inches(8.0), y + Inches(0.05), Inches(1.5), Inches(0.4),
            rate, size=17, color=BG_PRIMARY if is_total else SUCCESS, bold=True)
        txt(s, Inches(9.8), y + Inches(0.05), Inches(3), Inches(0.4),
            contrib, size=14, color=BG_PRIMARY if is_total else TEXT_MUTED)

    txt(s, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4),
        '最も効果が大きい工程: テスト（90〜95%短縮）— 開発と同時にAIがテストを自動生成',
        size=16, color=WARNING, bold=True)

    # ================================================================
    # 6: なぜ可能か — Claude Code Level 5
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    txt(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
        'なぜこれが可能なのか？', size=36, color=ACCENT, bold=True)
    txt(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
        '魔法ではなく「仕組み」— Claude Code の5段階最適化', size=18, color=TEXT_SECONDARY)

    levels = [
        ('1', '素のプロンプト', '毎回ルールを口頭で指示', ACCENT),
        ('2', '+ CLAUDE.md', 'ルールを自動読み込み', ACCENT),
        ('3', '+ Skills', '手順書をコマンドで注入', SUCCESS),
        ('4', '+ Hooks', 'フォーマット・テストを自動実行', SUCCESS),
        ('5', '+ Agents', 'セキュリティ・性能レビューを並行実行', WARNING),
    ]
    for i, (num, title, desc, color) in enumerate(levels):
        y = Inches(1.9 + i * 1.0)
        rect(s, Inches(0.8), y, Inches(11.5), Inches(0.85), BG_SURFACE)
        circle_num(s, Inches(1.1), y + Inches(0.0), num, color)
        txt(s, Inches(2.3), y + Inches(0.08), Inches(3.5), Inches(0.4),
            title, size=20, color=color, bold=True)
        txt(s, Inches(6.0), y + Inches(0.12), Inches(6), Inches(0.4),
            desc, size=16, color=TEXT_SECONDARY)

    txt(s, Inches(0.8), Inches(7.0), Inches(11), Inches(0.4),
        'Level 5 到達で、人間がやるのは「何を作るか決める」と「動作確認する」だけ',
        size=16, color=WARNING, bold=True)

    # ================================================================
    # 7: ユメログ紹介
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    txt(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
        'で、何を作ったの？', size=36, color=ACCENT, bold=True)
    txt(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
        '「夢を行動に変えるアプリ」— 3週間・1人・未経験で作った32機能',
        size=18, color=TEXT_SECONDARY)

    # 3ステップ
    steps = [
        ('1', '書き出す', '夢を言葉にする', ACCENT),
        ('2', '分解する', '目標→タスクに分解', SUCCESS),
        ('3', '続ける', '星座で成長を実感', WARNING),
    ]
    for i, (num, title, desc, color) in enumerate(steps):
        x = Inches(0.8 + i * 4.1)
        rect(s, x, Inches(1.8), Inches(3.8), Inches(1.8), BG_SURFACE)
        circle_num(s, x + Inches(0.2), Inches(2.0), num, color)
        txt(s, x + Inches(1.3), Inches(2.05), Inches(2.2), Inches(0.4),
            title, size=22, color=color, bold=True)
        txt(s, x + Inches(1.3), Inches(2.5), Inches(2.2), Inches(0.4),
            desc, size=15, color=TEXT_SECONDARY)

    features = [
        ('夢・目標・タスク管理', '階層で整理 + 発見ガイド'),
        ('活動予定', 'タイムライン + リマインド'),
        ('書籍管理', '本棚風UI + 読了レビュー'),
        ('統計・分析', 'グラフで活動を可視化'),
        ('星座システム', 'ゲーム感覚で継続'),
        ('ダッシュボード', 'カスタマイズ表示'),
    ]
    for i, (title, desc) in enumerate(features):
        col = i % 3
        row = i // 3
        x = Inches(0.8 + col * 4.1)
        y = Inches(4.0 + row * 1.6)
        rect(s, x, y, Inches(3.8), Inches(1.3), BG_SURFACE)
        txt(s, x + Inches(0.3), y + Inches(0.15), Inches(3.2), Inches(0.4),
            title, size=17, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        txt(s, x + Inches(0.3), y + Inches(0.6), Inches(3.2), Inches(0.5),
            desc, size=14, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)

    # ================================================================
    # 8: 開発者の想い
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    txt(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
        'なぜこのアプリを作ったのか', size=36, color=ACCENT, bold=True)

    rect(s, Inches(1.5), Inches(1.5), Inches(10.3), Inches(4.0), BG_SURFACE)
    multi(s, Inches(2.0), Inches(1.8), Inches(9.3), Inches(3.5), [
        {'text': '学生時代、夢を持てず将来に希望を見出せない時期がありました。',
         'size': 20, 'color': TEXT_SECONDARY, 'space': 24},
        {'text': '本との出会いを通じて', 'size': 20, 'color': TEXT_SECONDARY, 'space': 4},
        {'text': '「生き生きと生きるには夢が必要」', 'size': 28, 'color': ACCENT, 'bold': True, 'space': 4},
        {'text': 'と気づきました。', 'size': 20, 'color': TEXT_SECONDARY, 'space': 24},
        {'text': '夢を言語化し行動に移す仕組みを、アプリとして形にしました。',
         'size': 20, 'color': TEXT_SECONDARY, 'space': 8},
    ])

    txt(s, Inches(2), Inches(5.8), Inches(9), Inches(0.5),
        'ユメログは「最初の一歩」を支えるアプリです。',
        size=26, color=WARNING, bold=True, align=PP_ALIGN.CENTER)

    txt(s, Inches(2), Inches(6.5), Inches(9), Inches(0.4),
        'そしてこのアプリ自体が、AI駆動開発における私の「最初の一歩」でもあります。',
        size=16, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # ================================================================
    # 9: CTA
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    txt(s, Inches(1.5), Inches(0.6), Inches(10), Inches(0.8),
        '3週間で作ったアプリの品質、\n自分の目で確かめてみてください',
        size=34, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(1.5), Inches(2.0), Inches(10), Inches(0.5),
        'Webブラウザですぐに開けます。登録不要・完全無料。',
        size=18, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)

    if os.path.exists(qr):
        s.shapes.add_picture(qr, Inches(4.9), Inches(2.8), Inches(3.5), Inches(3.5))

    txt(s, Inches(2), Inches(6.5), Inches(9), Inches(0.4),
        'https://teppei19980914.github.io/GrowthEngine/',
        size=18, color=ACCENT, align=PP_ALIGN.CENTER)

    # ================================================================
    # 10: まとめ
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    txt(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
        'まとめ', size=36, color=ACCENT, bold=True)

    summaries = [
        ('未経験技術の学習コストがゼロに', SUCCESS),
        ('テスト工程が消滅（開発と同時に自動生成）', SUCCESS),
        ('開発期間が 6〜9分の1 に短縮', SUCCESS),
    ]
    for i, (text, color) in enumerate(summaries):
        y = Inches(1.5 + i * 1.4)
        rect(s, Inches(1.2), y, Inches(10.8), Inches(1.0), BG_SURFACE)
        circle_num(s, Inches(1.5), y + Inches(0.08), str(i + 1), color)
        txt(s, Inches(2.8), y + Inches(0.2), Inches(8.5), Inches(0.6),
            text, size=24, color=TEXT_PRIMARY, bold=True)

    rect(s, Inches(2), Inches(5.8), Inches(9.3), Inches(1.0), ACCENT,
         'AI駆動開発は「人間の潜在能力を引き出すエンジン」',
         size=22, color=BG_PRIMARY)

    txt(s, Inches(3.5), Inches(7.0), Inches(6), Inches(0.4),
        'ご清聴ありがとうございました',
        size=18, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # ================================================================
    out = os.path.join(os.path.dirname(__file__), 'LT_ユメログ.pptx')
    prs.save(out)
    print(f'Generated: {out}')


if __name__ == '__main__':
    create()
